import io
import subprocess
from pathlib import Path
from types import SimpleNamespace

import pymupdf
import pytest
from docx import Document
from flask import current_app
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader, PdfWriter
from werkzeug.datastructures import FileStorage

from app import create_app
from app.services.document_service import _safe_spreadsheet_value, office_to_pdf
from app.services.file_service import read_document_upload
from app.utils.errors import ToolError
from config import TestingConfig, config_by_name


def make_pdf(page_count=2, table=False):
    document = pymupdf.open()
    for index in range(page_count):
        page = document.new_page(width=300, height=220)
        page.insert_text((40, 35), f"Toolbox conversion page {index + 1}")
        if table:
            xs, ys = (40, 150, 260), (60, 100, 140)
            for x in xs:
                page.draw_line((x, ys[0]), (x, ys[-1]))
            for y in ys:
                page.draw_line((xs[0], y), (xs[-1], y))
            page.insert_text((50, 86), "Name")
            page.insert_text((160, 86), "Value")
            page.insert_text((50, 126), "Alpha")
            page.insert_text((160, 126), "42")
    data = document.tobytes()
    document.close()
    return io.BytesIO(data)


def test_preview_endpoint_renders_real_bounded_pages(client):
    response = client.post(
        "/api/pdf/preview",
        data={"pdf": (make_pdf(3), "source.pdf"), "limit": "2"},
        content_type="multipart/form-data",
        headers={"X-Requested-With": "Toolbox"},
    )
    assert response.status_code == 200
    payload = response.get_json()
    assert payload["page_count"] == 3
    assert payload["rendered_count"] == 2
    assert payload["truncated"] is True
    assert len(payload["pages"]) == 2
    assert payload["pages"][0]["data_url"].startswith("data:image/png;base64,iVBOR")
    assert payload["pages"][0]["width"] <= 180
    assert payload["pages"][0]["height"] <= 260
    assert response.headers["Cache-Control"].startswith("no-store")


def test_preview_rejects_non_pdf_and_obeys_server_cap(client, app):
    malformed = client.post(
        "/api/pdf/preview",
        data={"pdf": (io.BytesIO(b"not a pdf"), "bad.pdf")},
        content_type="multipart/form-data",
        headers={"X-Requested-With": "Toolbox"},
    )
    assert malformed.status_code == 400
    app.config["PDF_PREVIEW_MAX_PAGES"] = 1
    capped = client.post(
        "/api/pdf/preview",
        data={"pdf": (make_pdf(3), "source.pdf"), "limit": "999"},
        content_type="multipart/form-data",
    )
    assert capped.get_json()["rendered_count"] == 1


def test_pdf_to_jpg_and_png_routes(client):
    single = client.post(
        "/pdf-to-jpg/process",
        data={"pdf": (make_pdf(1), "source.pdf"), "pages": "1", "dpi": "96", "quality": "80"},
        content_type="multipart/form-data",
    )
    assert single.status_code == 200
    assert single.mimetype == "image/jpeg"
    assert single.data.startswith(b"\xff\xd8\xff")
    multiple = client.post(
        "/pdf-to-png/process",
        data={"pdf": (make_pdf(2), "source.pdf"), "pages": "all", "dpi": "72"},
        content_type="multipart/form-data",
    )
    assert multiple.status_code == 200
    assert multiple.mimetype == "application/zip"


def test_pdf_to_word_preserves_text(client):
    response = client.post(
        "/pdf-to-word/process",
        data={"pdf": (make_pdf(2), "source.pdf")},
        content_type="multipart/form-data",
    )
    assert response.status_code == 200
    converted = Document(io.BytesIO(response.data))
    text = "\n".join(paragraph.text for paragraph in converted.paragraphs)
    assert "Toolbox conversion page 1" in text
    assert "Toolbox conversion page 2" in text


def test_pdf_to_excel_extracts_detected_tables(client):
    response = client.post(
        "/pdf-to-excel/process",
        data={"pdf": (make_pdf(1, table=True), "table.pdf")},
        content_type="multipart/form-data",
    )
    assert response.status_code == 200, response.get_json() if response.is_json else response.data[:200]
    workbook = load_workbook(io.BytesIO(response.data))
    sheet = workbook[workbook.sheetnames[0]]
    values = [cell.value for row in sheet.iter_rows() for cell in row]
    assert "Name" in values
    assert "42" in values


def test_pdf_to_excel_does_not_fake_tables(client):
    response = client.post(
        "/pdf-to-excel/process",
        data={"pdf": (make_pdf(1), "plain.pdf")},
        content_type="multipart/form-data",
    )
    assert response.status_code == 400
    assert "no structured tables" in response.get_json()["error"].lower()


def test_pdf_to_powerpoint_creates_one_slide_per_page(client):
    response = client.post(
        "/pdf-to-powerpoint/process",
        data={"pdf": (make_pdf(3), "source.pdf")},
        content_type="multipart/form-data",
    )
    assert response.status_code == 200
    presentation = Presentation(io.BytesIO(response.data))
    assert len(presentation.slides) == 3
    assert all(len(slide.shapes) == 1 for slide in presentation.slides)


@pytest.mark.parametrize(
    ("route", "filename"),
    (("/jpg-to-pdf/process", "photo.jpg"), ("/png-to-pdf/process", "photo.png")),
)
def test_focused_image_to_pdf_routes(client, route, filename):
    from PIL import Image

    source = io.BytesIO()
    image_format = "JPEG" if filename.endswith(".jpg") else "PNG"
    Image.new("RGB", (40, 30), "navy").save(source, format=image_format)
    source.seek(0)
    response = client.post(
        route,
        data={"images": [(source, filename)], "page_size": "a4", "orientation": "portrait", "margin": "18", "fit": "contain"},
        content_type="multipart/form-data",
    )
    assert response.status_code == 200
    assert len(PdfReader(io.BytesIO(response.data)).pages) == 1


def test_spreadsheet_values_cannot_become_formulas():
    assert _safe_spreadsheet_value("=HYPERLINK(\"https://invalid\")").startswith("'")
    assert _safe_spreadsheet_value("@SUM(A1:A2)").startswith("'")
    assert _safe_spreadsheet_value("ordinary text") == "ordinary text"


def make_docx_upload(filename="source.docx"):
    document = Document()
    document.add_paragraph("Safe source document")
    data = io.BytesIO()
    document.save(data)
    data.seek(0)
    return FileStorage(
        stream=data,
        filename=filename,
        content_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    )


def test_office_validator_rejects_false_extensions(app):
    with app.app_context():
        fake = FileStorage(stream=io.BytesIO(b"not docx"), filename="fake.docx", content_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document")
        with pytest.raises(ToolError):
            read_document_upload(fake, 1024, {".docx"})
        wrong = make_docx_upload("source.docx")
        with pytest.raises(ToolError):
            read_document_upload(wrong, 100_000, {".xlsx"})


@pytest.mark.parametrize("route", ["/word-to-pdf/process", "/excel-to-pdf/process", "/powerpoint-to-pdf/process"])
def test_office_route_reports_missing_engine(client, monkeypatch, route):
    monkeypatch.setattr("app.services.document_service.find_libreoffice", lambda: None)
    response = client.post(
        route,
        data={"document": (io.BytesIO(b"placeholder"), "source.docx")},
        content_type="multipart/form-data",
    )
    assert response.status_code == 503
    assert "libreoffice" in response.get_json()["error"].lower()


def test_office_timeout_is_safe_and_cleans_workspace(app, monkeypatch):
    executable = Path(__import__("sys").executable)
    monkeypatch.setattr("app.services.document_service.find_libreoffice", lambda: str(executable))

    def timeout_run(command, **_kwargs):
        raise subprocess.TimeoutExpired(command, timeout=1)

    monkeypatch.setattr(subprocess, "run", timeout_run)
    with app.app_context():
        temp_root = Path(current_app.instance_path) / "tmp"
        before = set(temp_root.iterdir())
        with pytest.raises(ToolError) as caught:
            office_to_pdf(make_docx_upload(), "word-to-pdf")
        after = set(temp_root.iterdir())
    assert caught.value.status_code == 504
    assert before == after


def test_office_conversion_uses_safe_command_and_cleans_workspace(app, monkeypatch):
    captured = {}
    executable = Path(__import__("sys").executable)
    monkeypatch.setattr("app.services.document_service.find_libreoffice", lambda: str(executable))

    def fake_run(command, **kwargs):
        captured["command"] = command
        captured["kwargs"] = kwargs
        source = Path(command[-1])
        writer = PdfWriter()
        writer.add_blank_page(width=72, height=72)
        with (source.parent / f"{source.stem}.pdf").open("wb") as output:
            writer.write(output)
        return SimpleNamespace(returncode=0)

    monkeypatch.setattr(subprocess, "run", fake_run)
    with app.app_context():
        temp_root = Path(current_app.instance_path) / "tmp"
        before = set(temp_root.iterdir())
        result = office_to_pdf(make_docx_upload(), "word-to-pdf")
        after = set(temp_root.iterdir())
    assert result.startswith(b"%PDF")
    assert before == after
    assert captured["kwargs"]["shell"] is False
    assert captured["kwargs"]["timeout"] == app.config["DOCUMENT_CONVERSION_TIMEOUT_SECONDS"]
    assert "--headless" in captured["command"]
    assert "source.docx" not in captured["command"][-1]


def test_preview_rate_limit_returns_429(monkeypatch):
    class RateTestingConfig(TestingConfig):
        RATELIMIT_ENABLED = True
        RATELIMIT_STORAGE_URI = "memory://"
        RATELIMIT_PDF_PREVIEW = "2 per minute"

    monkeypatch.setitem(config_by_name, "rate-testing", RateTestingConfig)
    app = create_app("rate-testing")
    client = app.test_client()
    statuses = []
    for _index in range(3):
        statuses.append(client.post(
            "/api/pdf/preview",
            data={"pdf": (make_pdf(1), "source.pdf")},
            content_type="multipart/form-data",
            headers={"X-Requested-With": "Toolbox"},
        ).status_code)
    assert statuses == [200, 200, 429]
