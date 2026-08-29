"""Practical PDF/Office conversions with bounded temporary processing."""

import io
import os
import shutil
import subprocess
from pathlib import Path

import pymupdf
from docx import Document
from docx.shared import Inches
from flask import current_app
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches as PptxInches

from app.services.file_service import (
    generated_name,
    read_document_upload,
    temporary_workspace,
    validate_pdf_bytes,
)
from app.services.pdf_preview_service import open_pdf_document, render_page_png
from app.utils.errors import ToolError


OFFICE_INPUTS = {
    "word-to-pdf": {".doc", ".docx"},
    "excel-to-pdf": {".xls", ".xlsx"},
    "powerpoint-to-pdf": {".ppt", ".pptx"},
}


def _conversion_page_limit(document):
    limit = current_app.config["MAX_CONVERSION_PAGES"]
    if document.page_count > limit:
        raise ToolError(f"Convert PDFs with no more than {limit} pages at a time.")


def pdf_to_docx(upload):
    document, _data = open_pdf_document(upload)
    try:
        _conversion_page_limit(document)
        output_doc = Document()
        output_doc.core_properties.title = str(document.metadata.get("title") or "")[:255]
        section = output_doc.sections[0]
        section.top_margin = Inches(0.65)
        section.bottom_margin = Inches(0.65)
        section.left_margin = Inches(0.7)
        section.right_margin = Inches(0.7)
        for page_index, page in enumerate(document):
            blocks = sorted(page.get_text("dict").get("blocks", []), key=lambda item: (item.get("bbox", [0, 0])[1], item.get("bbox", [0])[0]))
            for block in blocks:
                if block.get("type") == 0:
                    lines = []
                    for line in block.get("lines", []):
                        text = "".join(span.get("text", "") for span in line.get("spans", [])).strip()
                        if text:
                            lines.append(text)
                    if lines:
                        output_doc.add_paragraph("\n".join(lines))
                elif block.get("type") == 1 and block.get("image"):
                    bbox = block.get("bbox", (0, 0, 360, 0))
                    width_inches = min(6.5, max(0.5, (bbox[2] - bbox[0]) / 72.0))
                    try:
                        output_doc.add_picture(io.BytesIO(block["image"]), width=Inches(width_inches))
                    except (ValueError, OSError):
                        current_app.logger.info("Skipped an unsupported embedded PDF image.")
            if page_index < document.page_count - 1:
                output_doc.add_page_break()
        output = io.BytesIO()
        output_doc.save(output)
        return output.getvalue()
    finally:
        document.close()


def _safe_spreadsheet_value(value):
    if isinstance(value, str) and value.lstrip().startswith(("=", "+", "-", "@")):
        return "'" + value
    return value


def pdf_to_xlsx(upload):
    document, _data = open_pdf_document(upload)
    try:
        _conversion_page_limit(document)
        workbook = Workbook()
        workbook.remove(workbook.active)
        table_number = 0
        for page_index, page in enumerate(document):
            try:
                tables = page.find_tables().tables
            except (AttributeError, RuntimeError, ValueError) as exc:
                current_app.logger.info("PDF table detection failed on page %s: %s", page_index + 1, exc)
                tables = []
            for table in tables:
                table_number += 1
                if table_number > 100:
                    raise ToolError("That PDF contains too many tables to convert safely.")
                worksheet = workbook.create_sheet(f"Page {page_index + 1} Table {table_number}"[:31])
                for row_index, row in enumerate(table.extract(), start=1):
                    for column_index, value in enumerate(row, start=1):
                        worksheet.cell(row_index, column_index, _safe_spreadsheet_value(value))
                worksheet.freeze_panes = "A2"
                for column in worksheet.columns:
                    values = [len(str(cell.value or "")) for cell in column]
                    worksheet.column_dimensions[column[0].column_letter].width = min(45, max(10, max(values, default=10) + 2))
        if not workbook.sheetnames:
            raise ToolError("No structured tables were detected in this PDF. Try a table-based PDF with visible rows and columns.")
        output = io.BytesIO()
        workbook.save(output)
        return output.getvalue()
    finally:
        document.close()


def pdf_to_pptx(upload):
    document, _data = open_pdf_document(upload)
    try:
        _conversion_page_limit(document)
        presentation = Presentation()
        first = document.load_page(0).rect
        slide_height = PptxInches(7.5)
        slide_width = int(slide_height * first.width / first.height)
        presentation.slide_width = slide_width
        presentation.slide_height = slide_height
        blank_layout = presentation.slide_layouts[6]
        total_pixels = 0
        for page in document:
            image, width, height = render_page_png(page, 120)
            total_pixels += width * height
            if total_pixels > current_app.config["MAX_CONVERSION_TOTAL_PIXELS"]:
                raise ToolError("That PDF is too large to convert safely at slide resolution.")
            slide = presentation.slides.add_slide(blank_layout)
            page_ratio = width / height
            slide_ratio = presentation.slide_width / presentation.slide_height
            if page_ratio > slide_ratio:
                placed_width = presentation.slide_width
                placed_height = int(placed_width / page_ratio)
                left, top = 0, int((presentation.slide_height - placed_height) / 2)
            else:
                placed_height = presentation.slide_height
                placed_width = int(placed_height * page_ratio)
                left, top = int((presentation.slide_width - placed_width) / 2), 0
            slide.shapes.add_picture(io.BytesIO(image), left, top, width=placed_width, height=placed_height)
        output = io.BytesIO()
        presentation.save(output)
        return output.getvalue()
    finally:
        document.close()


def find_libreoffice():
    configured = current_app.config.get("LIBREOFFICE_PATH", "")
    candidates = [configured] if configured else []
    candidates.extend(filter(None, [shutil.which("soffice.com"), shutil.which("soffice"), shutil.which("libreoffice")]))
    if os.name == "nt":
        for variable in ("ProgramFiles", "ProgramFiles(x86)"):
            base = os.environ.get(variable)
            if base:
                candidates.append(str(Path(base) / "LibreOffice" / "program" / "soffice.com"))
    for candidate in candidates:
        path = Path(candidate)
        if path.is_file():
            return str(path.resolve())
    return None


def libreoffice_available():
    return find_libreoffice() is not None


def office_to_pdf(upload, tool_slug):
    allowed = OFFICE_INPUTS.get(tool_slug)
    if not allowed:
        raise ToolError("That document conversion is unavailable.", status_code=404)
    executable = find_libreoffice()
    if not executable:
        raise ToolError(
            "Office-to-PDF conversion is unavailable on this server. Install LibreOffice or configure LIBREOFFICE_PATH.",
            status_code=503,
        )
    data, suffix = read_document_upload(
        upload,
        current_app.config["MAX_DOCUMENT_FILE_BYTES"],
        allowed,
    )
    with temporary_workspace() as workspace:
        source_path = workspace / generated_name(suffix)
        source_path.write_bytes(data)
        profile_path = workspace / "libreoffice-profile"
        profile_path.mkdir()
        command = [
            executable,
            "--headless",
            "--nologo",
            "--nodefault",
            "--nolockcheck",
            "--nofirststartwizard",
            f"-env:UserInstallation={profile_path.resolve().as_uri()}",
            "--convert-to",
            "pdf",
            "--outdir",
            str(workspace),
            str(source_path),
        ]
        try:
            completed = subprocess.run(
                command,
                cwd=str(workspace),
                capture_output=True,
                text=True,
                timeout=current_app.config["DOCUMENT_CONVERSION_TIMEOUT_SECONDS"],
                check=False,
                shell=False,
            )
        except subprocess.TimeoutExpired as exc:
            raise ToolError("Document conversion timed out. Try a smaller file.", status_code=504) from exc
        except OSError as exc:
            current_app.logger.warning("LibreOffice could not start: %s", exc)
            raise ToolError("The document conversion engine could not start.", status_code=503) from exc
        output_path = workspace / f"{source_path.stem}.pdf"
        if completed.returncode != 0 or not output_path.is_file():
            current_app.logger.warning("LibreOffice conversion failed with exit code %s", completed.returncode)
            raise ToolError("LibreOffice could not convert that document. Check that the file is valid and try again.")
        result = output_path.read_bytes()
        if len(result) > current_app.config["MAX_CONTENT_LENGTH"]:
            raise ToolError("The converted PDF is too large to return safely.")
        validate_pdf_bytes(result)
        return result
